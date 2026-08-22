"""Loader for PowerPoint (.ppt/.pptx) presentations.

Uses python-pptx directly (lighter than the `unstructured` package) to
extract text from every shape on every slide, one Document per slide so
slide-number metadata is preserved.
"""
import logging
import os
from typing import List

from pptx import Presentation
from langchain_core.documents import Document

logger = logging.getLogger(__name__)


def load_ppt(file_path: str) -> List[Document]:
    """Extract text from a PPT/PPTX file and return standardized Documents.

    Raises:
        FileNotFoundError: if the file does not exist.
        ValueError: if extraction fails or no slide contains extractable text.
    """
    if not os.path.isfile(file_path):
        raise FileNotFoundError(f"PPT file not found: {file_path}")

    try:
        presentation = Presentation(file_path)
    except Exception as exc:
        logger.error("Failed to open PPT %s: %s", file_path, exc)
        raise ValueError(f"Could not extract text from PPT: {file_path}") from exc

    filename = os.path.basename(file_path)
    docs = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        slide_text = "\n".join(texts).strip()
        if slide_text:
            docs.append(
                Document(
                    page_content=slide_text,
                    metadata={"source": filename, "type": "ppt", "slide": slide_number},
                )
            )

    if not docs:
        raise ValueError(
            f"No extractable text found in PPT (possibly image-only slides): {file_path}"
        )

    logger.info("Loaded PPT: %s (%d slide(s) with text)", filename, len(docs))
    return docs
