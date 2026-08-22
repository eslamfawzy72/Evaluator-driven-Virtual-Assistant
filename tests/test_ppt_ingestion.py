"""Verify the PPT loader works end-to-end through ingest() -> retrieve()."""
import os
import tempfile

from pptx import Presentation
from pptx.util import Inches

from ingestion.ingest import ingest
from rag.retriever import retrieve


def _build_test_pptx(path: str) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[1]  # title + content

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Ingestion Pipeline"
    slide1.placeholders[1].text = "Documents are chunked before being embedded."

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Caching"
    slide2.placeholders[1].text = "Redis caches retrieval results to avoid repeated searches."

    prs.save(path)


def test_ingest_ppt_and_retrieve():
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    _build_test_pptx(path)

    try:
        num_chunks = ingest(path, ".pptx")
        assert num_chunks > 0

        context = retrieve("What does Redis cache?", k=2)
        assert len(context) > 0
        assert any("retrieval results" in c["content"].lower() for c in context)
    finally:
        os.remove(path)
